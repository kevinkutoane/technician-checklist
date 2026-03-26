"""
make_presentation.py
Generates TechnicianChecklist-OpsHub.pptx using python-pptx.
Run: pip install python-pptx && python make_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colour palette ────────────────────────────────────────────────────────────
PRIMARY    = RGBColor(0x1A, 0x5F, 0xA8)   # #1A5FA8  brand blue
ACCENT     = RGBColor(0x0E, 0x8A, 0x6C)   # #0E8A6C  green
BG_DARK    = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B  dark sidebar
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF1, 0xF5, 0xF9)
MID_GREY   = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, BG_DARK)

    # Accent bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(3.2), Inches(10), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    add_textbox(slide, title_text,
                left=0.5, top=1.6, width=9, height=1.2,
                font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, subtitle_text,
                left=0.5, top=3.4, width=9, height=0.8,
                font_size=18, bold=False, color=MID_GREY, align=PP_ALIGN.CENTER)
    return slide


def add_content_slide(prs, title_text, body_lines,
                      title_color=PRIMARY, bg_color=WHITE,
                      text_color=TEXT_DARK, font_size=16):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, bg_color)

    # Title bar
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = title_color
    bar.line.fill.background()

    add_textbox(slide, title_text,
                left=0.3, top=0.1, width=9.4, height=0.8,
                font_size=24, bold=True, color=WHITE)

    # Body — multi-line text box starting below title bar
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(9.0), Inches(5.4))
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in body_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        if line.startswith("•") or line.startswith("▸"):
            p.space_before = Pt(4)
        if line.strip() == "":
            run.font.size = Pt(6)

    return slide


def add_two_col_slide(prs, title_text, left_lines, right_lines,
                      title_color=PRIMARY):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_bg(slide, WHITE)

    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = title_color
    bar.line.fill.background()

    add_textbox(slide, title_text,
                left=0.3, top=0.1, width=9.4, height=0.8,
                font_size=24, bold=True, color=WHITE)

    def _col(lines, left_in):
        txBox = slide.shapes.add_textbox(
            Inches(left_in), Inches(1.2), Inches(4.4), Inches(5.4))
        tf = txBox.text_frame; tf.word_wrap = True
        first = True
        for line in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run = p.add_run(); run.text = line
            run.font.size = Pt(15); run.font.color.rgb = TEXT_DARK
            if line.startswith("•"):
                p.space_before = Pt(4)

    _col(left_lines, 0.4)
    _col(right_lines, 5.2)

    # Divider
    div = slide.shapes.add_shape(
        1, Inches(5.0), Inches(1.2), Inches(0.02), Inches(5.4))
    div.fill.solid(); div.fill.fore_color.rgb = LIGHT_GREY
    div.line.fill.background()

    return slide


# ── Build presentation ────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── 1. Title ─────────────────────────────────────────────────────────────
    add_title_slide(prs,
        "GIBS Ops Hub",
        "Technician Checklist & Operations Management Platform  |  2025")

    # ── 2. Overview ──────────────────────────────────────────────────────────
    add_content_slide(prs, "System Overview", [
        "• Full-stack web application for AV/IT technicians and operations staff",
        "• Centralised platform replacing paper-based checklists and ad-hoc communication",
        "",
        "▸ Core value proposition:",
        "  – Daily classroom equipment checks with photo / note capture",
        "  – Asset agreement signing (digital signature)",
        "  – QA inspections with pass/fail grading",
        "  – Classroom handover records between shifts",
        "  – Hybrid classroom camera/mic/Teams tracking",
        "  – Week-ahead schedule from uploaded Excel file",
        "  – Admin panel with audit log, user management, analytics",
        "",
        "▸ Key design goals:",
        "  – Works offline (Service Worker PWA)",
        "  – Role-based access: technician vs admin",
        "  – Secure: session-based auth, CSRF-safe sameSite cookies, input sanitisation",
    ])

    # ── 3. User Roles ─────────────────────────────────────────────────────────
    add_two_col_slide(prs, "User Roles",
        left_lines=[
            "👷 Technician",
            "",
            "• Submit classroom checklists",
            "• Sign asset agreements",
            "• Complete QA checklists",
            "• Submit handover forms",
            "• View week-ahead schedule",
            "• View own dashboard / history",
            "• Equipment loan requests",
            "• Change personal preferences",
        ],
        right_lines=[
            "⚙️  Admin",
            "",
            "• All technician capabilities",
            "• Upload week-ahead XLSX",
            "• View & export full audit log",
            "• Manage users (create / disable)",
            "• View equipment status trends",
            "• Delete week-ahead uploads",
            "• View all submissions / history",
            "• Configure system settings",
        ]
    )

    # ── 4. Tech Stack ─────────────────────────────────────────────────────────
    add_two_col_slide(prs, "Technology Stack",
        left_lines=[
            "🖥️  Back-end",
            "",
            "• Node.js 18 + Express 4",
            "• SQLite via better-sqlite3",
            "• express-session + better-sqlite3",
            "  session store (rolling: true)",
            "• Multer — file uploads",
            "• SheetJS (xlsx) — XLSX parsing",
            "• Nodemailer — email alerts",
            "• node-cron — scheduled backup",
            "• bcrypt — password hashing",
        ],
        right_lines=[
            "🌐  Front-end & infra",
            "",
            "• Vanilla JS (no framework)",
            "• CSS custom properties (themes)",
            "• Service Worker — PWA / offline",
            "• Manifest.json — installable",
            "• localStorage — theme persistence",
            "• SignaturePad (CDN) — asset sign",
            "",
            "🧪  Testing",
            "• Jest + Supertest  |  243 tests",
            "• In-memory SQLite test DB",
        ]
    )

    # ── 5. Technician Features ────────────────────────────────────────────────
    add_content_slide(prs, "Technician Feature Map", [
        "Feature                 Route                    Key capability",
        "─────────────────────────────────────────────────────────────────────",
        "Classroom Checklist     /checklist               Equipment tick-list + notes, email alert on flag",
        "Asset Agreement         /onboarding              Digital sig capture, PDF download",
        "QA Checklist            /qa                      Pass/warn/fail grading, PDF export",
        "Handover Form           /handover                Issue log between shift handovers",
        "Week Ahead              /week-ahead              View schedule by day; admin uploads XLSX",
        "Dashboard               /dashboard               Summary cards, today's hybrid classrooms",
        "Equipment Loans         /loans                   Loan requests and return tracking",
        "Settings                /settings                Theme toggle (light/dark), password change",
    ], font_size=13)

    # ── 6. Dashboard ─────────────────────────────────────────────────────────
    add_content_slide(prs, "Dashboard", [
        "• Today's Hybrid Classrooms — card shows each room's setup status (camera/mic/Teams)",
        "• My Recent Activity — last 5 checklists and agreements",
        "• Week Ahead Summary card — event count for current week",
        "• Equipment Status Snapshot — flagged items",
        "",
        "Admin-only sections:",
        "• Equipment Status Trends — line chart of flags over 30 days (Chart.js)",
        "• Full submission history across all technicians",
        "• Audit log preview",
        "",
        "API endpoints:",
        "  GET /api/dashboard/today-progress",
        "  GET /api/dashboard/equipment-trends   (admin)",
    ])

    # ── 7. Classroom Checklist ────────────────────────────────────────────────
    add_content_slide(prs, "Classroom Checklist", [
        "• Select classroom → equipment list pre-loaded from admin config",
        "• Each item: ✅ OK  /  ⚠️ Issue  /  ❌ Fault",
        "• Notes field per item (pre-filled from last submission)",
        "• Flag triggers email alert to admin (nodemailer)",
        "",
        "DB tables:  classrooms, equipment, checklists, checklist_items",
        "",
        "Key routes:",
        "  GET  /api/checklists/classrooms",
        "  GET  /api/checklists/latest-notes?classroom_id=X",
        "  POST /api/checklists",
        "  GET  /api/checklists/history",
    ])

    # ── 8. Asset Agreement ────────────────────────────────────────────────────
    add_content_slide(prs, "Asset Agreement (Onboarding)", [
        "• Technician signs asset list on screen using SignaturePad",
        "• Signature stored as base-64 PNG in DB column signature_data",
        "• PDF download: server-side HTML → PDF via route GET /api/onboarding/:id/pdf",
        "",
        "DB table:  asset_agreements  (id, user_id, items JSON, signature_data, created_at)",
        "",
        "Key routes:",
        "  POST /api/onboarding",
        "  GET  /api/onboarding/history",
        "  GET  /api/onboarding/:id/pdf",
    ])

    # ── 9. QA Checklist ───────────────────────────────────────────────────────
    add_content_slide(prs, "QA Checklist", [
        "• Structured inspection: grouped checks (Audio, Video, Network, Safety…)",
        "• Each item: Pass / Warning / Fail",
        "• Overall grade calculated server-side",
        "• PDF export available from history",
        "",
        "DB table:  qa_submissions  (id, user_id, classroom_id, results JSON, grade, created_at)",
        "",
        "Key routes:",
        "  POST /api/qa",
        "  GET  /api/qa/history",
        "  GET  /api/qa/:id/pdf",
    ])

    # ── 10. Handover Form ─────────────────────────────────────────────────────
    add_content_slide(prs, "Classroom Handover", [
        "• Records classroom state at end of shift",
        "• Fields: classroom, outgoing tech, incoming tech, issues noted, status",
        "• Visible to all technicians and admin",
        "",
        "DB table:  handovers",
        "",
        "Key routes:",
        "  POST /api/handover",
        "  GET  /api/handover/history",
    ])

    # ── 11. Hybrid Classroom Notifications ───────────────────────────────────
    add_content_slide(prs, "Hybrid Classroom Setup Tracking", [
        "• Technicians mark camera, microphone, and Teams status per room",
        "• Dashboard card shows today's rooms requiring hybrid setup",
        "• Badge count on dashboard card updates dynamically",
        "",
        "API: GET /api/dashboard/today-progress",
        "",
        "Response shape:",
        '  { "date": "2025-02-24", "classrooms": [',
        '      { "id": 1, "name": "Room A", "camera": true,',
        '        "microphone": true, "teams": false }',
        '    ]',
        '  }',
    ], font_size=14)

    # ── 12. Week Ahead ────────────────────────────────────────────────────────
    add_content_slide(prs, "📅 Week Ahead (New Feature)", [
        "• Admin uploads a .xlsx file containing the week's event schedule",
        "• Parser extracts: date, time range, venue, company/course, pax, tech assigned",
        "• Results stored in week_ahead_events table; uploads tracked in week_ahead_uploads",
        "",
        "Client features:",
        "  – Day-tab navigation (Mon–Sun) with event count badges",
        "  – Week ◀ / ▶ navigation + jump-to-week",
        "  – After upload, view auto-jumps to the uploaded week",
        "  – Upload history table with delete (admin)",
        "",
        "Parser robustness:",
        "  – raw: false avoids Excel serial-date issues",
        "  – Scans merged cells for date header label",
        "  – Time-regex guard rejects summary rows (TOTAL DELEGATES, Day, Various…)",
        "  – Full + abbreviated month names supported (Jan–Dec / January–December)",
        "",
        "Key routes:  POST /api/week-ahead/upload  |  GET /api/week-ahead/events",
        "             GET /api/week-ahead/uploads  |  DELETE /api/week-ahead/uploads/:id",
    ], font_size=13)

    # ── 13. Admin Panel ───────────────────────────────────────────────────────
    add_content_slide(prs, "Admin Panel", [
        "Tabs:",
        "• Users          — create, list, disable/enable technician accounts",
        "• Classrooms     — add / edit rooms, assign equipment",
        "• Equipment      — asset register with category & serial",
        "• Audit Log      — timestamped log of all write operations",
        "• Week Ahead     — upload XLSX, view / delete upload history",
        "",
        "Audit log entries written on:",
        "  POST/PATCH/DELETE to checklist, onboarding, qa, handover, admin routes",
        "  Middleware helper: logAudit(req, action, detail)",
        "",
        "DB table:  audit_log (id, user_id, action, detail, ip, created_at)",
        "Route:     GET /api/audit-log  (admin only, paginated)",
    ])

    # ── 14. Security ──────────────────────────────────────────────────────────
    add_content_slide(prs, "Security Architecture", [
        "Authentication & session:",
        "  • bcrypt password hashing (cost 12)",
        "  • express-session with HttpOnly, SameSite=strict, Secure (prod)",
        "  • 8-hour session with rolling refresh on every request",
        "  • Session store persisted in SQLite (survives restart)",
        "",
        "Authorisation:",
        "  • requireAuth middleware  — blocks all unauthenticated API requests (401)",
        "  • requireAdmin middleware — blocks non-admin routes (403)",
        "  • Role column checked on every admin route",
        "",
        "Input safety:",
        "  • All user-generated content HTML-escaped with esc() before rendering",
        "  • Parameterised SQL queries via better-sqlite3 (no string interpolation)",
        "  • File upload: MIME + extension whitelist, 5 MB size limit",
        "  • Rate limiting on auth routes (express-rate-limit, skip in tests)",
        "",
        "Theme & data stored in localStorage only (non-sensitive; no session data)",
    ], font_size=14)

    # ── 15. API Endpoints ─────────────────────────────────────────────────────
    add_content_slide(prs, "API Endpoint Summary", [
        "Auth:        POST /api/auth/login  |  GET /api/auth/me  |  POST /logout",
        "Checklist:   POST /api/checklists  |  GET history  |  GET latest-notes",
        "Onboarding:  POST /api/onboarding  |  GET history  |  GET :id/pdf",
        "QA:          POST /api/qa          |  GET history  |  GET :id/pdf",
        "Handover:    POST /api/handover    |  GET history",
        "Dashboard:   GET  /api/dashboard/today-progress",
        "             GET  /api/dashboard/equipment-trends  (admin)",
        "Settings:    GET/PATCH /api/settings/preferences",
        "             PATCH     /api/settings/password",
        "Loans:       POST /api/loans  |  PATCH :id/return  |  GET history",
        "Admin:       GET/POST/PATCH /api/admin/users",
        "             GET/POST/PATCH /api/admin/classrooms",
        "             GET/POST/PATCH /api/admin/equipment",
        "             GET /api/audit-log",
        "Week Ahead:  POST /api/week-ahead/upload",
        "             GET  /api/week-ahead/events?week_start=YYYY-MM-DD",
        "             GET  /api/week-ahead/uploads",
        "             DELETE /api/week-ahead/uploads/:id  (admin)",
    ], font_size=12)

    # ── 16. Database Schema ───────────────────────────────────────────────────
    add_content_slide(prs, "Database Schema (SQLite)", [
        "users              — id, username, full_name, password_hash, role, active",
        "sessions           — sid, sess, expire  (session store)",
        "classrooms         — id, name, location, active",
        "equipment          — id, classroom_id, name, category, serial",
        "checklists         — id, user_id, classroom_id, submitted_at",
        "checklist_items    — id, checklist_id, equipment_id, status, notes",
        "asset_agreements   — id, user_id, items JSON, signature_data, created_at",
        "qa_submissions     — id, user_id, classroom_id, results JSON, grade, created_at",
        "handovers          — id, user_id, classroom_id, issues, status, created_at",
        "loans              — id, user_id, equipment_id, purpose, returned_at",
        "user_preferences   — user_id (PK), theme",
        "audit_log          — id, user_id, action, detail, ip, created_at",
        "week_ahead_uploads — id, user_id, filename, week_start, week_end, row_count, hash, created_at",
        "week_ahead_events  — id, upload_id, event_date, day_label, time_range, venue,",
        "                     company_course, contact_person, pax_campus, pax_zoom,",
        "                     lecturer, syndicates_other_venues, assigned_tech, it_support_required",
    ], font_size=12)

    # ── 17. Testing ───────────────────────────────────────────────────────────
    add_content_slide(prs, "Test Suite — 243 Tests", [
        "Framework: Jest + Supertest | SQLite in-memory test DB | --runInBand",
        "",
        "Test files:",
        "  tests/auth.test.js         — login, session, logout, me endpoint",
        "  tests/checklist.test.js    — submit, history, notes pre-fill, flagging",
        "  tests/onboarding.test.js   — signature, PDF generation, history",
        "  tests/qa.test.js           — grading logic, PDF, history",
        "  tests/dashboard.test.js    — today-progress, equipment trends, PDF history",
        "  tests/admin.test.js        — CRUD users/classrooms/equipment, audit log",
        "  tests/middleware.test.js   — requireAuth, requireAdmin, rate limiting",
        "  tests/weekahead.test.js    — XLSX upload, event retrieval, delete, access control",
        "",
        "Coverage thresholds (package.json):",
        "  lines ≥ 85%  |  branches ≥ 80%",
        "",
        "CI: GitHub Actions workflow on push to main",
    ])

    # ── 18. Recent Bug Fixes & Improvements ──────────────────────────────────
    add_content_slide(prs, "Recent Bug Fixes & Improvements", [
        "Bug Fix Phase 1:",
        "  ✅ Session rolling:true added — keeps sessions alive during active use",
        "  ✅ initNav: non-401 errors now show Retry overlay instead of blank page",
        "  ✅ Dashboard: silent catch on secondary API calls prevents white screen",
        "  ✅ CSS tablet breakpoint fixed (768px sidebar overlap)",
        "  ✅ Admin panel: null-check before accessing tab elements",
        "",
        "FOUC & Theme fixes:",
        "  ✅ localStorage theme cache — theme now applied before first paint",
        "  ✅ Inline script in <head> of all 11 pages reads opsHubTheme from localStorage",
        "  ✅ applyTheme() now persists theme to localStorage on every change",
        "",
        "Week Ahead parser fixes:",
        "  ✅ raw: false on sheet_to_json — prevents Excel serial-date values",
        "  ✅ Merged-cell scan for date header — searches all cells in DATE row",
        "  ✅ Time-regex guard — rejects TOTAL DELEGATES, Day, Various junk rows",
        "  ✅ Abbreviated month names (Jan, Feb … Dec) added to parser",
        "  ✅ Post-upload navigation — view jumps to the week that was just uploaded",
    ], font_size=13)

    # ── Save ──────────────────────────────────────────────────────────────────
    output = "TechnicianChecklist-OpsHub.pptx"
    prs.save(output)
    print(f"✅ Saved: {output}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    build()
